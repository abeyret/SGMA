"""Build a professional ECON 30 thesis presentation (8 slides)."""
from __future__ import annotations

import json
import re
import shutil
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "ECON30_Thesis_Presentation.pptx"
OUT_LEGACY = ROOT / "Copy of ECON107_ABEYRET.pptx"
BACKUP = ROOT / "Copy of ECON107_ABEYRET.backup.pptx"
WEBSITE = "https://econ301.vercel.app"

# Academic palette
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
CHARCOAL = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF4, 0xF3, 0xF0)
RULE = RGBColor(0xC8, 0xC0, 0xB4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_L = 0.85
MARGIN_R = 0.85
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


def load_county_well_data() -> list[tuple[str, int, int]]:
    data = (ROOT / "vercel_site" / "atlas_data.js").read_text(encoding="utf-8")
    m = re.search(r"const ATLAS = (\{.*\});", data, re.DOTALL)
    if not m:
        raise RuntimeError("Could not parse atlas_data.js")
    atlas = json.loads(m.group(1))
    rows = []
    for feat in atlas["counties"]["features"]:
        name = feat["properties"]["name"]
        pre = feat["properties"]["pre"]["well_failures_issue_start"]
        post = feat["properties"]["post"]["well_failures_issue_start"]
        rows.append((name, pre, post))
    rows.sort(key=lambda r: r[2], reverse=True)
    return rows


def make_well_failures_chart(path: Path) -> None:
    rows = load_county_well_data()
    counties = [r[0] for r in rows]
    pre = [r[1] for r in rows]
    post = [r[2] for r in rows]

    rcParams.update(
        {
            "font.family": "serif",
            "font.size": 11,
            "axes.spines.top": False,
            "axes.spines.right": False,
        }
    )
    fig, ax = plt.subplots(figsize=(8.5, 4.2), dpi=200)
    x = range(len(counties))
    w = 0.36
    ax.bar([i - w / 2 for i in x], pre, width=w, color="#1e3a5f", label="Pre-SGMA (2012–14)")
    ax.bar([i + w / 2 for i in x], post, width=w, color="#8b4513", label="Post-SGMA (2018–22)")
    ax.set_xticks(list(x))
    ax.set_xticklabels(counties, rotation=0)
    ax.set_ylabel("Dry-well reports (issue start date)")
    ax.set_title("Reported well failures rose sharply after SGMA", loc="left", fontsize=13, pad=12)
    ax.legend(frameon=False, loc="upper right")
    ax.yaxis.grid(True, linestyle="--", alpha=0.35)
    ax.set_axisbelow(True)
    fig.text(0.01, 0.01, "Source: DWR Household Water Supply Shortage Reporting System", fontsize=8, color="#666")
    fig.tight_layout()
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def make_mechanism_diagram(path: Path) -> None:
    """Publication-style aquifer cross-section."""
    w, h = 1800, 720
    img = Image.new("RGB", (w, h), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    try:
        title = ImageFont.truetype("times.ttf", 34)
        label = ImageFont.truetype("times.ttf", 24)
        small = ImageFont.truetype("times.ttf", 20)
    except OSError:
        title = ImageFont.load_default()
        label = title
        small = title

    draw.text((48, 28), "Mechanism: groundwater overdraft and land subsidence", fill=(30, 58, 95), font=title)

    top, bottom = 110, h - 70
    left, right = 120, w - 120
    mid = 430

    draw.rectangle([left, top, right, mid - 70], fill=(232, 220, 196))
    draw.rectangle([left, mid - 70, right, bottom], fill=(58, 110, 150))

    # water table
    draw.line([left, mid - 70, right, mid - 70], fill=(255, 255, 255), width=3)
    draw.line([left, mid + 35, right, mid + 35], fill=(255, 210, 120), width=3)
    draw.text((left + 20, mid - 108), "Original saturated zone", fill=(30, 58, 95), font=small)
    draw.text((left + 20, mid + 45), "Water table after sustained pumping", fill=(120, 80, 20), font=small)

    # land surface
    draw.line([left, top, right, top], fill=(40, 40, 40), width=4)
    draw.line([left, top + 42, right, top + 42], fill=(139, 69, 19), width=4)
    draw.text((right - 430, top + 2), "Original ground surface", fill=(40, 40, 40), font=small)
    draw.text((right - 470, top + 48), "Subsided ground surface", fill=(139, 69, 19), font=small)

    # well
    wx = right - 260
    draw.rectangle([wx, top + 42, wx + 18, mid + 35], fill=(90, 90, 90))
    draw.polygon([(wx + 9, top - 10), (wx - 35, top + 42), (wx + 53, top + 42)], fill=(30, 58, 95))
    draw.text((wx + 30, top + 55), "Extraction", fill=(30, 58, 95), font=label)

    for x in range(left + 80, wx - 40, 95):
        draw.line([x, mid - 70, x, mid + 35], fill=(255, 255, 255), width=2)
        draw.polygon([(x - 8, mid + 45), (x + 8, mid + 45), (x, mid + 58)], fill=(255, 255, 255))

    draw.text((left + 180, mid + 95), "Loss of pore-water pressure compacts aquifer sediments", fill=(255, 255, 255), font=label)

    steps = [
        "Pumping exceeds natural recharge",
        "Hydraulic head declines",
        "Effective stress increases; pores collapse",
        "Land surface subsides",
    ]
    sx = left
    for i, step in enumerate(steps, 1):
        draw.ellipse([sx, bottom + 8, sx + 28, bottom + 36], fill=(30, 58, 95))
        draw.text((sx + 9, bottom + 12), str(i), fill=(255, 255, 255), font=small)
        draw.text((sx + 36, bottom + 12), step, fill=(40, 40, 40), font=small)
        sx += 420

    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rule(slide, top: float) -> None:
    line = slide.shapes.add_shape(1, Inches(MARGIN_L), Inches(top), Inches(CONTENT_W), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_header(slide, title: str, *, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(0.55), Inches(CONTENT_W), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Georgia"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Calibri"
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(4)
    add_rule(slide, 1.35)


def add_bullets(
    slide,
    items: list[str],
    *,
    left=MARGIN_L,
    top=1.65,
    width=6.0,
    size=18,
    spacing=6,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = CHARCOAL
        p.space_after = Pt(spacing)
        p.level = 0


def add_figure(
    slide,
    path: Path,
    left: float,
    top: float,
    width: float,
    *,
    height: float | None = None,
    caption: str | None = None,
) -> None:
    if height is None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
        cap_top = top + 3.05
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))
        cap_top = top + height + 0.08
    if caption:
        cap = slide.shapes.add_textbox(Inches(left), Inches(cap_top), Inches(width + 0.2), Inches(0.45))
        p = cap.text_frame.paragraphs[0]
        p.text = caption
        p.font.name = "Calibri"
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.font.italic = True


def add_footer(slide, text: str, *, slide_num: int | None = None, total: int = 8) -> None:
    box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(6.95), Inches(CONTENT_W - 0.5), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    if slide_num is not None:
        num = slide.shapes.add_textbox(Inches(SLIDE_W - 1.0), Inches(6.95), Inches(0.5), Inches(0.35))
        np = num.text_frame.paragraphs[0]
        np.text = str(slide_num)
        np.font.name = "Calibri"
        np.font.size = Pt(10)
        np.font.color.rgb = GRAY
        np.alignment = PP_ALIGN.RIGHT


def add_metric_block(slide, metrics: list[tuple[str, str, str]], *, left=MARGIN_L, top=1.75, width=6.2) -> None:
    """Headline value, label, source — stacked cleanly."""
    y = top
    for value, label, source in metrics:
        vb = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(width), Inches(0.55))
        vp = vb.text_frame.paragraphs[0]
        vp.text = value
        vp.font.name = "Georgia"
        vp.font.size = Pt(28)
        vp.font.bold = True
        vp.font.color.rgb = NAVY

        lb = slide.shapes.add_textbox(Inches(left), Inches(y + 0.48), Inches(width), Inches(0.45))
        lp = lb.text_frame.paragraphs[0]
        lp.text = label
        lp.font.name = "Calibri"
        lp.font.size = Pt(16)
        lp.font.color.rgb = CHARCOAL

        sb = slide.shapes.add_textbox(Inches(left), Inches(y + 0.88), Inches(width), Inches(0.3))
        sp = sb.text_frame.paragraphs[0]
        sp.text = source
        sp.font.name = "Calibri"
        sp.font.size = Pt(10)
        sp.font.color.rgb = GRAY
        sp.font.italic = True
        y += 1.35


def build() -> None:
    if OUT.exists() and not BACKUP.exists():
        shutil.copy2(OUT, BACKUP)

    chart_path = ASSETS / "thesis_well_failures_chart.png"
    diagram_path = ASSETS / "thesis_mechanism_diagram.png"
    make_well_failures_chart(chart_path)
    make_mechanism_diagram(diagram_path)

    prs = new_prs()
    blank = prs.slide_layouts[6]

    # Slide 1 — Motivation
    s1 = prs.slides.add_slide(blank)
    set_bg(s1)
    add_header(s1, "The San Joaquin Central Valley is sinking")
    add_metric_block(
        s1,
        [
            (">30 cm/yr", "Peak subsidence rates in parts of the valley", "Faunt et al. (2016)"),
            ("14 km³", "Subsidence volume, 2006–2022", "Knight & Lee (2024)"),
            ("9 m", "Historical subsidence near benchmark S661, 1925–1977", "USGS / DWR"),
        ],
        width=6.4,
    )
    add_bullets(
        s1,
        [
            "Groundwater extraction has exceeded recharge for decades across California’s most productive agricultural region.",
            "Land subsidence is a measurable physical externality of overdraft—not merely a hydrologic abstraction.",
        ],
        top=5.05,
        width=6.4,
        size=16,
    )
    add_figure(
        s1,
        ASSETS / "intro_slide2.png",
        left=7.35,
        top=1.55,
        width=5.1,
        height=3.05,
        caption="USGS subsidence benchmarks, San Joaquin Valley (1925–1977; 1988–2016).",
    )
    add_footer(s1, "Alexandra Beyret · ECON 30 · San Joaquin Valley Groundwater & SGMA Equity", slide_num=1)

    # Slide 2 — Subsidence repair costs
    s2 = prs.slides.add_slide(blank)
    set_bg(s2)
    add_header(s2, "Economic costs of subsidence and infrastructure repair")
    add_metric_block(
        s2,
        [
            ("$889 million", "Additional federal funding for San Joaquin Valley canal repairs", "March 2026 appropriation"),
            ("$1.87 billion", "Estimated aggregate housing value at risk from subsidence-related flooding", "July 2025 study"),
            ("$10k–$30k", "Typical household well repair or replacement costs in dry-well reports", "DWR dry-well database"),
        ],
        width=6.3,
    )
    add_bullets(
        s2,
        [
            "Subsidence damages rigid infrastructure—canals, bridges, levees, and conveyance systems require continuous maintenance.",
            "Repair costs are a direct transfer from groundwater overdraft to public agencies, irrigators, and households.",
        ],
        top=5.05,
        width=6.3,
        size=16,
    )
    add_figure(
        s2,
        ASSETS / "intro_slide1.png",
        left=7.35,
        top=1.55,
        width=5.1,
        height=3.05,
        caption="Canal embankment repair along a Central Valley conveyance corridor.",
    )
    add_footer(s2, "Infrastructure costs illustrate the scale of physical damages from land subsidence.", slide_num=2)

    # Slide 3 — Mechanism
    s3 = prs.slides.add_slide(blank)
    set_bg(s3)
    add_header(s3, "How overdraft translates into subsidence", subtitle="A compacted-aquifer mechanism")
    add_figure(s3, diagram_path, left=0.75, top=1.55, width=11.85, height=4.55)
    add_footer(
        s3,
        "When pumping lowers hydraulic head, effective stress on aquifer grains rises and fine-grained sediments compact irreversibly.",
        slide_num=3,
    )

    # Slide 4 — Overdraft consequences
    s4 = prs.slides.add_slide(blank)
    set_bg(s4)
    add_header(s4, "Overdraft imposes multiple economic externalities")
    col1 = [
        "Dry and failing wells — shallow domestic wells fail first; rural households bear outage costs.",
        "Water quality degradation — deeper pumping can mobilize salinity and legacy contaminants.",
        "Conveyance losses — subsidence reduces gravity-flow capacity in canals and aqueducts.",
    ]
    col2 = [
        "Agricultural adjustment — fallowed acreage and shifts in cropping intensity.",
        "Farm consolidation — net loss of 4,783 small farms (under 180 ac) across eight counties, 2012–2022.",
        "Uneven burden — disadvantaged communities face higher well-failure exposure (CalEnviroScreen × dry wells).",
    ]
    add_bullets(s4, col1, left=MARGIN_L, top=1.65, width=5.8, size=17)
    add_bullets(s4, col2, left=6.95, top=1.65, width=5.8, size=17)
    add_footer(s4, "Source: NASS farm operations; DWR dry-well reporting; project county summary.", slide_num=4)

    # Slide 5 — Management costs
    s5 = prs.slides.add_slide(blank)
    set_bg(s5)
    add_header(s5, "Groundwater management under SGMA", subtitle="Policy costs and implementation burden since 2014")
    add_metric_block(
        s5,
        [
            ("45 plans", "Groundwater Sustainability Plans in the San Joaquin Valley basin", "DWR GSP registry"),
            ("24 approved", "14 under review · 6 post state intervention · 1 incomplete", "Status as of project build"),
            ("2014", "Sustainable Groundwater Management Act — first statewide groundwater regulation", "California Water Code"),
        ],
        width=6.2,
    )
    add_bullets(
        s5,
        [
            "Compliance requires monitoring networks, demand reduction, fallowing, and potentially costly infrastructure retrofits.",
            "Implementation is spatially uneven: approved plans coexist with basins still under state oversight.",
            "Central policy question: who pays for sustainability—large irrigators, small farmers, or domestic well users?",
        ],
        top=5.05,
        width=6.2,
        size=16,
    )
    add_figure(
        s5,
        ASSETS / "intro_slide3.png",
        left=7.35,
        top=1.55,
        width=5.1,
        height=3.05,
        caption="State and federal conveyance infrastructure that SGMA is designed to protect.",
    )
    add_footer(s5, "SGMA shifts groundwater from open-access extraction toward regulated sustainability.", slide_num=5)

    # Slide 6 — Why here, why now
    s6 = prs.slides.add_slide(blank)
    set_bg(s6)
    add_header(s6, "Why the Central Valley—and why now?")
    add_bullets(
        s6,
        [
            "Scale: eight counties supply a disproportionate share of U.S. agricultural output and depend heavily on groundwater.",
            "Timing: SGMA (2014) coincided with drought, declining surface allocations, and intensified pumping.",
            "Evidence of stress: reported well failures rose from 182 (pre-SGMA) to 1,853 (post-SGMA) across the study counties.",
            "Policy urgency: without sustainable yield, subsidence and rural water insecurity will continue to impose rising social costs.",
        ],
        top=1.65,
        width=5.9,
        size=16,
    )
    add_figure(
        s6,
        chart_path,
        left=7.0,
        top=1.65,
        width=5.5,
        height=3.35,
        caption="Dry-well reports by county, binned by approximate issue start date.",
    )
    add_footer(s6, "The valley combines hydrologic vulnerability, economic dependence on irrigation, and new regulatory institutions.", slide_num=6)

    # Slide 7 — Research question
    s7 = prs.slides.add_slide(blank)
    set_bg(s7, LIGHT_GRAY)
    add_header(s7, "Research question")
    qbox = s7.shapes.add_textbox(Inches(1.1), Inches(1.85), Inches(11.1), Inches(2.4))
    qtf = qbox.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.text = (
        "Will SGMA-induced groundwater regulation stop the consequences of overextraction, "
        "or shift the costs onto small farmers and disadvantaged communities?"
    )
    qp.font.name = "Georgia"
    qp.font.size = Pt(26)
    qp.font.color.rgb = NAVY
    qp.font.italic = True
    qp.alignment = PP_ALIGN.CENTER

    add_bullets(
        s7,
        [
            "Empirical approach: link pre-existing harm (well failures), physical costs (subsidence / drawdown), policy adjustment (fallowing), and distributional outcomes (farm size, CalEnviroScreen).",
            "County-level panel (n = 8) with pre/post SGMA comparisons; scatter views test mechanism stories in an interactive atlas.",
            "Descriptive regressions (HC1 robust SE): fallow change strongly tracks groundwater change (R² = 0.95); equity mechanisms remain heterogeneous across counties.",
        ],
        top=4.35,
        width=11.0,
        size=17,
    )
    add_footer(s7, "Goal: evaluate whether sustainability policy mitigates overdraft externalities or reallocates them.", slide_num=7)

    # Slide 8 — Website
    s8 = prs.slides.add_slide(blank)
    set_bg(s8)
    add_header(s8, "Interactive data atlas", subtitle="SGMA Equity Pathways — San Joaquin Valley")
    add_bullets(
        s8,
        [
            "Scroll narrative maps subsidence, groundwater, GSP status, and equity metrics across eight counties.",
            "Atlas tools: county/GSP cards, mechanism scatter plots, and choropleth comparisons (Pre, Post, Δ).",
            "Data: CASGEM groundwater, DWR dry-well reporting, CDL fallow acreage, NASS farm size, CalEnviroScreen.",
        ],
        top=1.75,
        width=7.2,
        size=18,
    )
    url_box = s8.shapes.add_textbox(Inches(1.0), Inches(4.85), Inches(11.3), Inches(0.9))
    up = url_box.text_frame.paragraphs[0]
    up.text = WEBSITE
    up.font.name = "Calibri"
    up.font.size = Pt(28)
    up.font.bold = True
    up.font.color.rgb = NAVY
    up.alignment = PP_ALIGN.CENTER
    add_footer(s8, "Alexandra Beyret · ECON 30 · Questions and feedback welcome", slide_num=8)

    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")
    try:
        shutil.copy2(OUT, OUT_LEGACY)
        print(f"Also updated {OUT_LEGACY}")
    except OSError as exc:
        print(f"Could not overwrite {OUT_LEGACY} (close it in PowerPoint): {exc}")


if __name__ == "__main__":
    build()
